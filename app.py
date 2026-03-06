import streamlit as st
import pandas as pd
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches
import io

# 1. KONFIGURASI HALAMAN (Full Screen & Judul Tab)
st.set_page_config(
    page_title="Sales Analytics Pro - Papa",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="collapsed" # Sidebar otomatis tertutup saat dibuka
)

# 2. CUSTOM CSS (Full Dark Mode, Font Raksasa, & Layout Rata)
st.markdown("""
    <style>
    /* Paksa Tema Gelap Total */
    .stApp {
        background-color: #0E1117 !important;
        color: #FFFFFF !important;
    }

    /* JUDUL UTAMA - 80px & Rata Tengah */
    h1 {
        font-size: 80px !important;
        font-weight: 900 !important;
        color: #60A5FA !important;
        text-align: center !important;
        margin-bottom: 10px !important;
        line-height: 1.1 !important;
    }

    /* SUB-JUDUL - 40px & Rata Tengah */
    h2 {
        font-size: 40px !important;
        font-weight: 400 !important;
        color: #CBD5E1 !important;
        text-align: center !important;
        margin-top: 0px !important;
        margin-bottom: 50px !important;
        border-bottom: none !important;
    }

    /* SIDEBAR - Panel Kontrol Besar */
    [data-testid="stSidebar"] {
        background-color: #1E293B !important;
        min-width: 400px !important;
    }
    [data-testid="stSidebar"] h2 {
        font-size: 35px !important;
        font-weight: bold !important;
        color: #F8FAFC !important;
    }

    /* KARTU INSTRUKSI - Simetris & Rata */
    .instruction-container {
        display: flex;
        justify-content: center;
        gap: 20px;
        margin-top: 20px;
    }
    .instruction-card {
        background-color: #1E293B;
        padding: 30px;
        border-radius: 20px;
        border: 3px solid #3B82F6;
        text-align: center;
        flex: 1; /* Biar lebar sama rata */
        display: flex;
        flex-direction: column;
        align-items: center;
        min-height: 350px; /* Biar tinggi sama rata */
        box-shadow: 0 10px 20px rgba(0,0,0,0.5);
    }
    .instruction-card h3 {
        font-size: 30px !important;
        margin-top: 15px !important;
        color: #F8FAFC !important;
    }
    .instruction-card p {
        font-size: 22px !important;
        color: #94A3B8 !important;
    }

    /* RESPONSIVE MOBILE */
    @media (max-width: 768px) {
        h1 { font-size: 45px !important; }
        h2 { font-size: 24px !important; }
        .instruction-container { flex-direction: column; }
        .instruction-card { min-height: auto; }
    }

    /* Memperbesar Elemen Lain */
    .stButton>button {
        font-size: 22px !important;
        height: 3.5em !important;
        font-weight: bold !important;
    }
    [data-testid="stMetricValue"] {
        font-size: 50px !important;
    }
    </style>
    """, unsafe_allow_html=True)

# --- SIDEBAR (PANEL KONTROL) ---
with st.sidebar:
    st.markdown("## 📥 Panel Kontrol")
    uploaded_file = st.file_uploader("Upload File Excel/CSV Papa", type=['xlsx', 'csv'])
    st.markdown("---")
    st.markdown("### 🛠️ Status")
    if uploaded_file:
        st.success("✅ File Berhasil Dimasukkan")
    else:
        st.info("💡 Menunggu Upload...")
    st.caption("Klik tanda panah (>) di pojok kiri atas untuk menutup kembali menu ini.")

# --- LOGIKA DASHBOARD ---
if uploaded_file:
    try:
        # 1. Baca Data
        if uploaded_file.name.endswith('.csv'):
            df_raw = pd.read_csv(uploaded_file, header=None)
        else:
            excel = pd.ExcelFile(uploaded_file)
            sheet = st.sidebar.selectbox("Pilih Sheet:", excel.sheet_names) if len(excel.sheet_names) > 1 else excel.sheet_names[0]
            df_raw = pd.read_excel(uploaded_file, sheet_name=sheet, header=None)

        # 2. Bersihkan Data (Format Papa)
        df_raw = df_raw.dropna(how='all', axis=0).dropna(how='all', axis=1).reset_index(drop=True)
        weeks = df_raw.iloc[0].ffill()
        prods = df_raw.iloc[1]
        data_body = df_raw.iloc[2:].copy()
        
        headers = []
        for w, p in zip(weeks[1:], prods[1:]):
            headers.append(f"{str(w).replace('nan','')} - {str(p).replace('nan','')}".strip(" -"))
            
        df_temp = pd.DataFrame(data_body.iloc[:, 1:].values, columns=headers)
        df_temp['Metrik'] = data_body.iloc[:, 0].values
        
        # Filter baris yang ada angkanya (Buang judul kota/sales)
        mask = df_temp.drop('Metrik', axis=1).apply(lambda r: pd.to_numeric(r, errors='coerce').notnull().any(), axis=1)
        df_temp = df_temp[mask].reset_index(drop=True)
        
        # Melt & Pivot
        df_melted = df_temp.melt(id_vars=['Metrik'], var_name='Kategori', value_name='Nilai')
        df_final = df_melted.pivot_table(index='Kategori', columns='Metrik', values='Nilai', aggfunc='first').reset_index()
        
        for c in df_final.columns:
            if c != 'Kategori': df_final[c] = pd.to_numeric(df_final[c], errors='coerce').fillna(0)

        # --- TAMPILAN DASHBOARD ---
        st.markdown(f"<h1>Laporan: {uploaded_file.name}</h1>", unsafe_allow_html=True)
        
        m1, m2 = st.columns(2)
        m1.metric("Data Minggu", f"{len(df_final)} Kolom")
        m2.metric("Jumlah Item", f"{len(df_final.columns)-1} Baris")

        st.markdown("---")

        pilihan = st.selectbox("🎯 Pilih Data Penjualan (Font Besar):", [c for c in df_final.columns if c != 'Kategori'])
        
        # Render Grafik
        fig = px.bar(df_final, x='Kategori', y=pilihan, text_auto='.2s', 
                     color_discrete_sequence=['#60A5FA'], title=f"Analisis {pilihan}")
        
        fig.update_layout(template="plotly_dark", font=dict(size=20), height=650)
        st.plotly_chart(fig, use_container_width=True)

        # Download PPT
        st.markdown("### 💾 Opsi Laporan")
        col_ppt, col_tab = st.columns(2)
        
        with col_ppt:
            def buat_ppt(grafik, judul):
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.title.text = judul
                img_io = io.BytesIO(grafik.to_image(format="png", width=1200, height=700))
                slide.shapes.add_picture(img_io, Inches(0.5), Inches(1.5), width=Inches(9))
                out = io.BytesIO()
                prs.save(out)
                return out.getvalue()

            if st.button("🚀 Buat Slide PowerPoint"):
                data_ppt = buat_ppt(fig, f"Laporan {pilihan}")
                st.download_button("📥 Download PPT Sekarang", data_ppt, f"Laporan_{pilihan}.pptx")

        with col_tab:
            with st.expander("🔍 Lihat Detail Tabel Data"):
                st.dataframe(df_final, use_container_width=True)

    except Exception as e:
        st.error(f"Gagal memproses data: {e}")

else:
    # --- TAMPILAN AWAL (LANDING PAGE) ---
    st.markdown("<h1>Portal Analisis Data Anda</h1>", unsafe_allow_html=True)
    st.markdown("<h2>Dashboard eksekutif untuk monitoring laporan mingguan secara real-time.</h2>", unsafe_allow_html=True)
    
    
    
    st.markdown("""
    <div class="instruction-container">
        <div class="instruction-card">
            <h1 style='font-size: 70px; margin:0;'>📂</h1>
            <h3>1. Unggah</h3>
            <p>Buka panel kontrol di kiri atas, lalu masukkan file Excel Papa.</p>
        </div>
        <div class="instruction-card">
            <h1 style='font-size: 70px; margin:0;'>📊</h1>
            <h3>2. Pantau</h3>
            <p>Lihat tren penjualan otomatis dengan grafik yang besar dan jelas.</p>
        </div>
        <div class="instruction-card">
            <h1 style='font-size: 70px; margin:0;'>🎞️</h1>
            <h3>3. Presentasi</h3>
            <p>Simpan hasil ke PowerPoint untuk langsung dipresentasikan.</p>
        </div>
    </div>
    """, unsafe_allow_html=True)
